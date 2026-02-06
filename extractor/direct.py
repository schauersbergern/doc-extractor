"""Modus 1: Direkte Textextraktion aus PPTX-XML.

Schnell, kein GPU, exakt für regulären Text.
Limitierung: Kein Text aus eingebetteten Bildern, kein Layout-Kontext.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.table import Table

from .models import SlideData, TableData, Timer
from .utils import estimate_tokens

logger = logging.getLogger(__name__)


def _extract_table(table: Table) -> TableData:
    rows_data = []
    for row in table.rows:
        row_text = []
        for cell in row.cells:
            cell_text = "\n".join(
                p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()
            )
            row_text.append(cell_text)
        rows_data.append(row_text)

    if len(rows_data) > 1:
        return TableData(headers=rows_data[0], rows=rows_data[1:])
    elif rows_data:
        return TableData(headers=[], rows=rows_data)
    return TableData(headers=[], rows=[])


def _extract_shape_text(shape: BaseShape) -> tuple[list[str], list[TableData]]:
    texts = []
    tables = []

    # Gruppierte Shapes rekursiv
    if hasattr(shape, "shapes"):
        try:
            for child in shape.shapes:
                ct, ctb = _extract_shape_text(child)
                texts.extend(ct)
                tables.extend(ctb)
            return texts, tables
        except Exception:
            pass

    if shape.has_table:
        tables.append(_extract_table(shape.table))
        return texts, tables

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)

    return texts, tables


def extract_direct(
    pptx_path: str | Path,
    slide_numbers: list[int] | None = None,
    include_notes: bool = False,
) -> list[SlideData]:
    """Direkte Textextraktion aus PPTX."""
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"Nicht gefunden: {pptx_path}")

    prs = Presentation(str(pptx_path))
    results = []

    for idx, slide in enumerate(prs.slides, start=1):
        if slide_numbers and idx not in slide_numbers:
            continue

        with Timer() as timer:
            slide_data = SlideData(slide_number=idx, extraction_method="direct")

            if slide.shapes.title:
                slide_data.title = slide.shapes.title.text.strip()

            all_texts = []
            all_tables = []

            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes.title:
                    continue
                texts, tables = _extract_shape_text(shape)
                all_texts.extend(texts)
                all_tables.extend(tables)

            slide_data.content = "\n".join(all_texts)
            slide_data.tables = all_tables

            if include_notes and slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_data.notes = notes_frame.text.strip()

        slide_data.extraction_time_seconds = timer.elapsed
        slide_data.token_count = estimate_tokens(slide_data.content)

        results.append(slide_data)
        logger.info(
            f"Slide {idx}: {len(all_texts)} Blöcke, "
            f"{len(all_tables)} Tabellen, {timer.elapsed:.3f}s"
        )

    return results
